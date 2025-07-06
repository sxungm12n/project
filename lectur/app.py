from flask import Flask, request, jsonify, render_template, send_from_directory, session
from openai import AzureOpenAI
from werkzeug.utils import secure_filename
import os
from datetime import datetime, timedelta
from lecture_recorder import LectureRecorder
from dotenv import load_dotenv
from flask_session import Session
import re
import requests
from docx import Document
from pptx import Presentation
import openpyxl
import PyPDF2
import io
from urllib.parse import quote

# .env 파일 로드
load_dotenv()

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.secret_key = 'your-secret-key'  # 세션을 위한 비밀 키
app.config['SESSION_TYPE'] = 'filesystem'  # 세션을 파일 시스템에 저장
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(days=7)  # 세션 유효 기간을 7일로 설정

# 업로드 폴더가 없으면 생성
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

# Azure OpenAI 설정
endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT")
api_key = os.getenv("AZURE_OPENAI_API_KEY")

# Azure Cognitive Search 설정
search_endpoint = os.getenv("AZURE_SEARCH_ENDPOINT")
search_index = os.getenv("AZURE_SEARCH_INDEX")
search_key = os.getenv("AZURE_SEARCH_KEY")
semantic_config = os.getenv("AZURE_SEARCH_SEMANTIC_CONFIG")

# OpenAI 클라이언트 생성
client = AzureOpenAI(
    azure_endpoint=endpoint,
    api_key=api_key,
    api_version="2024-05-01-preview"
)

# LectureRecorder 인스턴스 생성
recorder = LectureRecorder()

# 세션 초기화
Session(app)

# 이전 질문과 답변을 저장할 변수
last_question = None
last_answer = None

# Microsoft Translator API 설정
translator_endpoint = os.getenv("TRANSLATOR_ENDPOINT")
translator_key = os.getenv("TRANSLATOR_KEY")

@app.route('/')
def index():
    # 세션 초기화
    session.clear()
    return render_template('index.html')

@app.route('/main')
def main():
    # 세션 초기화
    session.clear()
    return render_template('index.html')

@app.route('/set_summary', methods=['POST'])
def set_summary():
    try:
        data = request.json
        summary = data.get('summary')
        
        if not summary:
            return jsonify({"success": False, "error": "앗, 깜빡하셨나요? 요약 내용을 살짝만 입력해주시면 제가 멋지게 정리해드릴게요! 💡"})
        
        # 세션에 요약 저장 (임시로만 저장)
        session['last_summary'] = summary
        session.permanent = False  # 세션을 임시로만 설정
        
        return jsonify({"success": True, "message": "정말 멋져요! 요약이 무사히 저장됐어요. 궁금한 게 있으면 언제든 또 불러주세요! 🌟"})
        
    except Exception as e:
        return jsonify({"success": False, "error": f"앗, 잠깐 문제가 생겼어요! 너무 걱정 마세요. 잠시 후 다시 시도해주시면 금방 해결해드릴게요 😅\n(에러: {str(e)})"})

@app.route('/chat', methods=['POST'])
def chat():
    global last_question, last_answer
    try:
        data = request.json
        user_message = data['messages'][0]['content']
        
        # 마지막 질문 저장
        last_question = user_message

        # 기본 시스템 메시지 설정
        base_system_message = {
            "role": "system",
            "content": """당신은 따뜻하고 열정적인 과학 선생님이자 친구입니다! 🌟

학생들의 질문에 다음과 같은 방식으로 답변해주세요:

1. **항상 따뜻하고 격려하는 톤**을 유지하세요 - "정말 좋은 질문이에요!", "궁금해하는 모습이 멋져요!" 같은 표현을 자주 사용하세요

2. **복잡한 개념은 일상적인 예시**로 설명하세요 - 학생이 쉽게 이해할 수 있도록 친숙한 상황에 비유해주세요

3. **단계별로 차근차근** 설명하세요 - 어려운 내용은 작은 단계로 나누어 설명해주세요

4. **과학적 호기심을 자극**하는 재미있는 사실이나 추가 정보를 포함하세요 - "재미있는 점은..." 같은 표현을 사용하세요

5. **적절한 이모지**를 사용하여 친근감을 더하세요 - 😊 🌟 💡 🔬 ⚡ 등

6. **칭찬과 격려**를 자주 해주세요 - "정말 잘 이해하고 있네요!", "이런 질문을 하는 것 자체가 대단해요!"

7. **다음 단계나 추가 질문**을 자연스럽게 유도하세요 - "더 궁금한 게 있으면 언제든 물어보세요!"

8. **실수나 어려움**에 대해서는 위로와 함께 해결책을 제시하세요 - "괜찮아요! 함께 해결해봐요!"

항상 학생의 관점에서 생각하고, 과학을 재미있고 흥미롭게 만들어주세요! 🚀"""
        }

        # 강의 내용 관련 키워드 체크
        lecture_related_keywords = ["강의", "수업", "강의내용", "수업내용", "강의 내용", "수업 내용"]
        is_lecture_query = any(keyword in user_message for keyword in lecture_related_keywords)

        # 강의 데이터가 있는 경우
        if 'last_summary' in session and session['last_summary']:
            try:
                # 강의 내용 기반 답변 시도
                lecture_system_message = {
                    "role": "system",
                    "content": f"""{base_system_message['content']}
                    
                    최근 강의 내용:
                    {session['last_summary']}
                    
                    위 강의 내용을 기반으로 답변하되, 강의 내용과 관련이 없는 질문은 일반적인 과학 지식을 바탕으로 답변해주세요."""
                }
                
                messages = [lecture_system_message] + data['messages']
                lecture_response = client.chat.completions.create(
                    model=deployment,
                    messages=messages,
                    temperature=0.9,  # 창의성 증가
                    max_tokens=2000,  # 최대 토큰 수 증가
                    top_p=0.95,      # 다양성 증가
                    frequency_penalty=0.5,  # 반복 감소
                    presence_penalty=0.5    # 새로운 토픽 선호
                )
                
                lecture_answer = lecture_response.choices[0].message.content
                if not "죄송합니다" in lecture_answer and not "관련된 내용을 찾을 수 없습니다" in lecture_answer:
                    last_answer = lecture_answer
                    return jsonify({
                        "response": lecture_answer,
                        "source": "lecture"
                    })
            except Exception as e:
                print(f"Lecture-based response error: {str(e)}")
        # 강의 데이터가 없는 경우 강의 관련 질문이면 안내 메시지 반환
        elif is_lecture_query:
            return jsonify({
                "response": "앗, 아직 강의 내용이 없어요! 🤔\n\n괜찮아요! 먼저 강의 파일을 업로드해주시면, 제가 열심히 공부해서 더 정확한 답변을 드릴 수 있어요! 📚✨\n\n강의 파일은 상단의 업로드 버튼을 통해 쉽게 업로드하실 수 있어요. 함께 해봐요! 💪",
                "source": "system"
            })

        # RAG 데이터 기반 답변 시도
        try:
            rag_messages = [base_system_message] + data['messages']
            rag_response = client.chat.completions.create(
                model=deployment,
                messages=rag_messages,
                temperature=0.9,      # 창의성 증가
                max_tokens=2000,      # 최대 토큰 수 증가
                top_p=0.95,          # 다양성 증가
                frequency_penalty=0.5,  # 반복 감소
                presence_penalty=0.5,   # 새로운 토픽 선호
                extra_body={
                    "data_sources": [
                        {
                            "type": "azure_search",
                            "parameters": {
                                "endpoint": search_endpoint,
                                "index_name": search_index,
                                "semantic_configuration": semantic_config,
                                "query_type": "semantic",
                                "top_n_documents": 10,  # 검색 문서 수 증가
                                "authentication": {
                                    "type": "api_key",
                                    "key": search_key
                                }
                            }
                        }
                    ]
                }
            )
            
            rag_answer = rag_response.choices[0].message.content
            if not "관련 정보를 찾을 수 없습니다" in rag_answer:
                # RAG 데이터의 출처 정보 추가
                citations = []
                if 'data_source_results' in rag_response.choices[0].message:
                    for source in rag_response.choices[0].message.data_source_results:
                        if 'title' in source:
                            citations.append(source['title'])
                
                citation_text = ""
                if citations:
                    citation_text = "\n\n---\n참조: " + ", ".join(citations)
                
                last_answer = rag_answer + citation_text
                return jsonify({
                    "response": rag_answer + citation_text,
                    "source": "rag"
                })
        except Exception as e:
            print(f"RAG response error: {str(e)}")

        # 일반 GPT 답변
        try:
            gpt_messages = [base_system_message] + data['messages']
            gpt_response = client.chat.completions.create(
                model=deployment,
                messages=gpt_messages,
                temperature=0.9,      # 창의성 증가
                max_tokens=2000,      # 최대 토큰 수 증가
                top_p=0.95,          # 다양성 증가
                frequency_penalty=0.5,  # 반복 감소
                presence_penalty=0.5    # 새로운 토픽 선호
            )
            
            gpt_answer = gpt_response.choices[0].message.content
            last_answer = gpt_answer
            return jsonify({
                "response": gpt_answer,
                "source": "gpt"
            })
            
        except Exception as e:
            error_message = str(e)
            if "DeploymentNotFound" in error_message:
                return jsonify({
                    "response": "앗, 지금 AI 서비스에 잠깐 문제가 생겼어요! 😅\n\n너무 걱정 마세요. 잠시 후에 다시 물어봐주시면 제가 더 열심히 답변해드릴게요! 🙏\n\n잠깐만 기다려주세요! ⏳",
                    "source": "system"
                })
            else:
                return jsonify({
                    "response": "음... 생각하는 중에 살짝 문제가 생겼어요! 🤔\n\n괜찮아요! 다시 한 번 질문해주시면 제가 더 잘 답변해드릴 수 있을 거예요! 💪\n\n함께 해봐요! 😊",
                    "source": "system"
                })
        
    except Exception as e:
        return jsonify({
            "response": "앗, 지금 시스템에 작은 문제가 생겼어요! 😅\n\n너무 걱정 마세요. 잠시 후에 다시 시도해주시면 금방 해결해드릴게요! 🙏\n\n잠깐만 기다려주세요! ⏳",
            "source": "system"
        })

@app.route('/summarize_last', methods=['GET'])
def summarize_last():
    global last_answer
    if not last_answer:
        return jsonify({
            "error": "앗, 아직 질문을 하지 않으셨네요! 😊\n\n괜찮아요! 먼저 궁금한 점을 물어봐주시면, 제가 멋지게 요약해드릴게요! 💡\n\n함께 시작해봐요! 🚀",
            "source": "system"
        })
    
    try:
        # 답변을 요약하도록 시스템 메시지 설정
        system_message = {
            "role": "system",
            "content": "다음 답변을 한 문장으로 간단히 요약해주세요."
        }
        
        # 답변을 요약
        response = client.chat.completions.create(
            model=deployment,
            messages=[
                system_message,
                {"role": "user", "content": last_answer}
            ],
            temperature=0.7,
            max_tokens=150
        )
        
        summary = response.choices[0].message.content.strip()
        
        return jsonify({
            "summary": summary,
            "source": "summary"
        })
        
    except Exception as e:
        error_message = str(e)
        if "DeploymentNotFound" in error_message:
            return jsonify({
                "error": "앗, 지금 요약 기능에 잠깐 문제가 생겼어요! 😅\n\n너무 걱정 마세요. 잠시 후에 다시 시도해주시면 제가 더 열심히 요약해드릴게요! 🙏",
                "source": "system"
            })
        elif "quota" in error_message.lower() or "limit" in error_message.lower():
            return jsonify({
                "error": "앗, 지금 AI 서비스가 너무 바빠서 잠깐 쉬어야 할 것 같아요! 😅\n\n괜찮아요! 잠시 후에 다시 시도해주시면 금방 도와드릴게요! 💪",
                "source": "system"
            })
        else:
            return jsonify({
                "error": "앗, 요약하는 중에 살짝 문제가 생겼어요! 😅\n\n괜찮아요! 다시 한 번 시도해주시면 제가 더 잘 요약해드릴 수 있을 거예요! 💪",
                "source": "system"
            })

@app.route('/search', methods=['POST'])
def search():
    global last_answer
    try:
        if not last_answer:
            return jsonify({
                "error": "앗, 아직 대화 내용이 없어요! 😊\n\n괜찮아요! 먼저 궁금한 점을 물어봐주시면, 제가 멋진 검색 결과를 보여드릴게요! 🔍\n\n함께 시작해봐요! 🚀",
                "source": "system"
            })

        # 마지막 답변에서 키워드 추출
        keyword_response = client.chat.completions.create(
            model=deployment,
            messages=[
                {
                    "role": "system",
                    "content": "다음 답변에서 가장 중요한 단일 키워드를 추출해주세요. 예: '마그마', '인공지능', '삼성전자'"
                },
                {"role": "user", "content": f"다음 답변에서 가장 중요한 단일 키워드를 추출해주세요:\n\n{last_answer}"}
            ],
            temperature=0.3,
            max_tokens=50
        )
        keyword = keyword_response.choices[0].message.content.strip()
        
        # 키워드 정제
        keyword = keyword.replace('-', '').strip()  # 하이픈 제거
        if keyword.startswith('"') or keyword.startswith("'"):
            keyword = keyword[1:]
        if keyword.endswith('"') or keyword.startswith("'"):
            keyword = keyword[:-1]
        keyword = keyword.strip()
        
        print(f"Extracted keyword: {keyword}")
        
        if not keyword:
            return jsonify({
                "error": "앗, 키워드를 추출하는 데 살짝 문제가 생겼어요! 😅\n\n괜찮아요! 다시 한 번 질문해주시면 제가 더 잘 찾아드릴 수 있을 거예요! 💪",
                "source": "system"
            })
        
        # 위키백과 검색 URL 생성
        encoded_keyword = quote(keyword)
        wiki_url = f"https://ko.wikipedia.org/wiki/{encoded_keyword}"
        
        # 메시지 생성 (마크다운 형식)
        message = f"""마지막 답변에서 추출한 키워드 **{keyword}**에 대한 위키백과 검색 결과입니다:

🔍 [위키백과에서 {keyword} 검색하기]({wiki_url})"""
        
        return jsonify({
            "response": message,
            "source": "search"
        })
        
    except Exception as e:
        print(f"Error in search: {str(e)}")
        return jsonify({
            "error": f"앗, 검색 중에 잠깐 문제가 생겼어요! 😅\n\n너무 걱정 마세요. 잠시 후에 다시 시도해주시면 금방 도와드릴게요! 🙏\n\n(에러: {str(e)})",
            "source": "system"
        })

def extract_keyword(question):
    """질문에서 핵심 키워드를 추출하는 함수"""
    
    # 불용어 목록 확장
    stopwords = [
        # 의문사
        '무엇', '뭐', '어떻', '어떤', '어떻게', '왜', '언제', '누구', '얼마나', '몇',
        # 조사
        '이', '가', '은', '는', '을', '를', '의', '에', '에서', '로', '으로', '와', '과', '이나', '나',
        # 대명사
        '나', '너', '우리', '저희', '그것', '이것', '저것', '그', '이', '저',
        # 접속사
        '그리고', '또는', '또', '및', '혹은',
        # 부사
        '매우', '너무', '아주', '잘', '더', '덜', '많이',
        # 동사/형용사 어미
        '이다', '하다', '되다', '있다', '없다', '이란', '하나요', '인가요', '일까요', '인지',
        # 기타 불용어
        '알려줘', '설명해줘', '가르쳐줘', '말해줘', '보여줘', '해줘', '좀', '제발', '부탁'
    ]
    
    # 질문에서 특수문자, 숫자 제거 및 소문자 변환
    cleaned = re.sub(r'[^\w\s]', ' ', question.lower())
    cleaned = re.sub(r'\d+', ' ', cleaned)
    
    # 단어 분리
    words = cleaned.split()
    
    # 불용어 제거 및 2글자 이상인 단어만 선택
    keywords = [word for word in words if word not in stopwords and len(word) >= 2]
    
    # 중복 제거 및 빈도수 기준 정렬
    from collections import Counter
    word_counts = Counter(keywords)
    
    # 가장 빈도수 높은 단어 선택 (없으면 None 반환)
    if word_counts:
        most_common = word_counts.most_common(1)[0][0]
        return most_common
    
    return None

def extract_text(file_path):
    """Extract text from various file types"""
    print(f"Extracting text from: {file_path}")
    try:
        file_ext = os.path.splitext(file_path)[1].lower()
        print(f"File extension: {file_ext}")
        
        if file_ext == '.wav':
            # WAV 파일 처리
            print("Processing WAV file...")
            recorder = LectureRecorder()
            text = recorder.request_STT(file_path)
            print(f"STT result: {text[:100]}...")  # 처음 100자만 출력
            if text.startswith('⚠️'):
                print(f"STT error: {text}")
                return text
            return text
            
        elif file_ext == '.txt':
            # 텍스트 파일 처리
            print("Processing text file...")
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
                
        elif file_ext in ['.docx', '.pptx', '.xlsx', '.xls', '.pdf']:
            # Office 문서 및 PDF 처리
            print(f"Processing {file_ext} file...")
            if file_ext == '.docx':
                doc = Document(file_path)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            elif file_ext == '.pptx':
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return '\n'.join(text)
            elif file_ext in ['.xlsx', '.xls']:
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text = []
                for sheet in wb:
                    for row in sheet.iter_rows():
                        text.append(' '.join([str(cell.value) for cell in row if cell.value]))
                return '\n'.join(text)
            elif file_ext == '.pdf':
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = []
                    for page in reader.pages:
                        text.append(page.extract_text())
                    return '\n'.join(text)
                    
        else:
            error_msg = f"⚠️ 지원되지 않는 파일 형식입니다: {file_ext}"
            print(error_msg)
            return error_msg
            
    except Exception as e:
        error_msg = f"⚠️ 텍스트 추출 중 오류가 발생했습니다: {str(e)}"
        print(f"Error in extract_text: {error_msg}")
        return error_msg

def allowed_file(filename):
    """Check if the file extension is allowed"""
    print(f"Checking file: {filename}")
    
    if not filename or '.' not in filename:
        print("No filename or no extension")
        return False
        
    ext = filename.rsplit('.', 1)[1].lower()
    print(f"File extension: {ext}")
    
    ALLOWED_EXTENSIONS = {'wav', 'txt', 'docx', 'pptx', 'xlsx', 'xls', 'pdf'}
    is_allowed = ext in ALLOWED_EXTENSIONS
    print(f"Allowed extensions: {ALLOWED_EXTENSIONS}")
    print(f"Is extension allowed: {is_allowed}")
    
    return is_allowed

@app.route('/upload', methods=['POST'])
def upload_file():
    print("Upload request received")
    
    if 'file' not in request.files:
        print("No file in request")
        return jsonify({'error': '앗, 깜빡하셨나요? 파일을 올려주시면 제가 바로 도와드릴게요! 🚀'}), 400
    
    file = request.files['file']
    if file.filename == '':
        print("Empty filename")
        return jsonify({'error': '파일 이름이 비어있어요! 실수할 수도 있죠 😊 다시 한 번 파일을 선택해주실래요?'}), 400
    
    print(f"Received file: {file.filename}")
    
    # 파일 확장자 검사
    if not allowed_file(file.filename):
        print(f"Invalid file extension: {file.filename}")
        return jsonify({'error': '이런! 이 파일은 아직 제가 공부를 못했어요. 혹시 .wav, .txt, .docx, .pptx, .xlsx, .xls, .pdf 파일로 다시 한 번 도전해주실래요? 💡'}), 400
    
    try:
        # 파일 저장 (원본 파일명 유지)
        filename = file.filename
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        print(f"Saving file to: {file_path}")
        file.save(file_path)
        
        # 텍스트 추출
        print("Starting text extraction...")
        extracted_text = extract_text(file_path)
        print(f"Extracted text: {extracted_text[:100]}...")  # 처음 100자만 출력
        
        if extracted_text.startswith("⚠️"):
            print(f"Error in text extraction: {extracted_text}")
            return jsonify({'error': f'파일에서 텍스트를 꺼내는 데 살짝 문제가 생겼어요. 파일 형식이나 내용을 한 번만 더 확인해주실래요?\n{extracted_text}\n그래도 잘 안되면 언제든 저를 불러주세요! 🙏'}), 400
        
        # 요약 생성
        print("Generating summary...")
        summary = generate_summary(extracted_text)
        
        # 세션에 요약 저장
        session['last_summary'] = summary
        
        # 파일 삭제
        print("Cleaning up...")
        os.remove(file_path)
        
        return jsonify({
            'message': '정말 멋져요! 파일이 무사히 도착했어요. 이제 결과를 확인해볼까요? 🌟',
            'text': extracted_text,
            'summary': summary
        })
        
    except Exception as e:
        print(f"Error in upload_file: {str(e)}")
        return jsonify({'error': f'앗, 업로드 과정에서 잠깐 문제가 생겼어요! 너무 걱정 마시고, 잠시 후 다시 시도해주시면 금방 도와드릴게요! 😅\n(에러: {str(e)})'}), 500

def generate_summary(text):
    try:
        response = client.chat.completions.create(
            model=deployment,
            messages=[
                {
                    "role": "system",
                    "content": "너는 텍스트를 받으면 그 내용을 분석하고 요약해주는 AI야. 또한 출력을 할 때는 markdown 형식으로 출력을 해야 해."
                },
                {
                    "role": "user",
                    "content": text
                }
            ],
            temperature=0.7,
            top_p=0.95,
            max_tokens=1200
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"요약 생성 중 오류가 발생했습니다: {str(e)}"

@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory('output', filename, as_attachment=True)

@app.route('/tts', methods=['POST'])
def text_to_speech():
    try:
        data = request.get_json()
        text = data.get('text', '')
        
        if not text:
            return jsonify({'error': '앗, 변환할 텍스트가 비어있어요! 한 줄만 입력해주시면 바로 멋진 음성으로 바꿔드릴게요 🎤'}), 400
            
        # TTS API 호출
        endpoint = os.getenv("TTS_ENDPOINT")
        headers = {
            "Content-Type": "application/ssml+xml",
            "Ocp-Apim-Subscription-Key": os.getenv("TTS_KEY"),
            "X-Microsoft-OutputFormat": "riff-24khz-16bit-mono-pcm"
        }
        
        # HTML 태그 제거
        text = re.sub(r'<[^>]+>', '', text)
        
        body = f"""
        <speak version='1.0' xml:lang='ko-KR'>
            <voice name='ko-KR-SeoHyeonNeural'>
                <prosody rate="150%">
                    {text}
                </prosody>
            </voice>
        </speak>
        """
        
        response = requests.post(endpoint, headers=headers, data=body)
        
        if response.status_code == 200:
            # 오디오 파일 저장
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"tts_{timestamp}.wav"
            filepath = os.path.join('static', 'audio', filename)
            
            # audio 디렉토리가 없으면 생성
            os.makedirs(os.path.join('static', 'audio'), exist_ok=True)
            
            with open(filepath, 'wb') as audio_file:
                audio_file.write(response.content)
                
            # 파일 URL 반환
            audio_url = f'/static/audio/{filename}'
            return jsonify({'audio_url': audio_url, 'message': '와우! 음성 파일이 완성됐어요! 아래에서 직접 들어보세요 🎵'})
        else:
            print(f"TTS API Error: {response.status_code} - {response.text}")
            return jsonify({'error': f'앗, TTS 변환에 실패했어요. 잠시 후 다시 시도해주시면 제가 더 열심히 해볼게요! (상태 코드: {response.status_code})'}), 500
            
    except Exception as e:
        print(f"TTS Error: {str(e)}")
        return jsonify({'error': f'음성 변환 중 잠깐 문제가 생겼어요! 너무 걱정 마시고, 다시 한 번만 시도해주실래요? 😅\n(에러: {str(e)})'}), 500

@app.route('/translate', methods=['POST'])
def translate():
    try:
        data = request.get_json()
        text = data.get('text', '')
        target_lang = data.get('target_lang', 'en')
        
        if not text:
            return jsonify({'error': '번역할 내용을 입력해주시면, 원하는 언어로 뚝딱 바꿔드릴게요! ✨'}), 400
            
        # Microsoft Translator API 호출
        headers = {
            'Ocp-Apim-Subscription-Key': os.getenv("TRANSLATOR_KEY"),
            'Ocp-Apim-Subscription-Region': 'eastus',
            'Content-Type': 'application/json'
        }
        
        params = {
            'api-version': '3.0',
            'to': target_lang
        }
        
        body = [{'text': text}]
        
        response = requests.post(
            os.getenv("TRANSLATOR_ENDPOINT"),
            headers=headers,
            params=params,
            json=body
        )
        
        if response.status_code == 200:
            result = response.json()
            if result and len(result) > 0 and 'translations' in result[0]:
                translated_text = result[0]['translations'][0]['text']
                return jsonify({'translated': translated_text, 'message': '번역이 뚝딱! 완료됐어요. 궁금한 게 있으면 언제든 또 물어봐주세요 😊'})
            else:
                print(f"Translation API response error: {response.text}")
                return jsonify({'error': '앗, 번역 결과를 처리하는 데 문제가 생겼어요. 잠시 후 다시 시도해주실래요?'}), 500
        elif response.status_code == 401:
            print(f"Translation API authentication error: {response.text}")
            return jsonify({'error': '앗, 번역 서비스 인증에 실패했어요. 혹시 관리자 선생님께 문의해주시면 빠르게 도와드릴 수 있어요!'}), 401
        else:
            print(f"Translation API error: {response.status_code} - {response.text}")
            return jsonify({'error': f'번역 API에서 잠깐 문제가 생겼어요. 너무 걱정 마시고, 잠시 후 다시 시도해주실래요? (상태 코드: {response.status_code})'}), response.status_code
            
    except Exception as e:
        print(f"Translation error: {str(e)}")
        return jsonify({'error': f'번역 중 잠깐 문제가 생겼어요! 다시 한 번만 시도해주시면 금방 도와드릴게요 😅\n(에러: {str(e)})'}), 500

@app.route('/generate_image', methods=['POST'])
def generate_image():
    global last_question, last_answer
    try:
        if not last_answer:
            return jsonify({
                "error": "아직 대화 내용이 없어요! 😊\n\n괜찮아요! 먼저 궁금한 점을 물어봐주시면, 제가 멋진 이미지를 만들어드릴 준비가 되어 있답니다 😊",
                "source": "system"
            })

        # 마지막 답변에서 키워드 추출
        keyword_response = client.chat.completions.create(
            model=deployment,
            messages=[
                {
                    "role": "system",
                    "content": "답변에서 가장 중요한 과학 관련 키워드 하나만 추출해주세요."
                },
                {"role": "user", "content": f"다음 답변에서 가장 중요한 과학 관련 키워드를 추출해주세요:\n\n{last_answer}"}
            ],
            temperature=0.3,
            max_tokens=50
        )
        keyword = keyword_response.choices[0].message.content.strip()

        # 프롬프트 자동 생성
        prompt = f"중학생 과학 교과서에 나오는 스타일로 {keyword} 그려줘."

        # DALL-E API 호출
        endpoint = os.getenv("DALLE_ENDPOINT")
        headers = {
            "Content-Type": "application/json",
            "api-key": os.getenv("DALLE_KEY")
        }
        body = {
            "prompt": prompt,
            "n": 1,
            "size": "1024x1024"
        }

        response = requests.post(endpoint, headers=headers, json=body)

        if response.status_code == 200:
            response_json = response.json()
            image_url = response_json["data"][0]["url"]
            return jsonify({
                "image_url": image_url,
                "keyword": keyword,
                "source": "image",
                "message": "두근두근! 이미지를 완성했어요. 아래에서 확인해보시고, 또 궁금한 게 있으면 언제든 말씀해 주세요! 🖼️"
            })
        else:
            print(f"DALL-E API Error: {response.status_code} - {response.text}")
            return jsonify({
                "error": f"앗, 이미지 생성에 실패했어요. 너무 걱정 마시고, 잠시 후 다시 시도해주시면 제가 더 열심히 해볼게요! (상태 코드: {response.status_code})",
                "source": "system"
            }), 500

    except Exception as e:
        print(f"Image generation error: {str(e)}")
        return jsonify({
            "error": f"이미지 생성 중 잠깐 문제가 생겼어요! 다시 한 번만 시도해주시면 금방 도와드릴게요 😅\n(에러: {str(e)})",
            "source": "system"
        }), 500

@app.route('/upload_message', methods=['POST'])
def upload_message():
    try:
        data = request.json
        message = data.get('message')
        
        if not message:
            return jsonify({"success": False, "error": "앗, 메시지가 비어있어요! 한 줄만 입력해주시면 바로 저장해드릴게요 😊"})
        
        # 메시지를 세션에 저장
        session['last_message'] = message
        session.permanent = False  # 세션을 임시로만 설정
        
        return jsonify({
            "success": True,
            "message": "메시지가 무사히 업로드됐어요! 궁금한 게 있으면 언제든 편하게 물어봐주세요 😊"
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": f"앗, 메시지 업로드 중 잠깐 문제가 생겼어요! 너무 걱정 마시고, 잠시 후 다시 시도해주시면 금방 도와드릴게요 😅\n(에러: {str(e)})"})

if __name__ == '__main__':
    # 로컬에서만 실행 (127.0.0.1)
    app.run(host='127.0.0.1', port=5500, debug=True) 